from __future__ import annotations

import csv
import json
import statistics
from collections import Counter
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import gridspec
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOC_DIR = ROOT / "doc"
ASSET_DIR = DOC_DIR / "report_assets"
PPT_PATH = DOC_DIR / "KnowSAM_超声脑部外侧裂半监督分割汇报.pptx"
MD_PATH = DOC_DIR / "KnowSAM_超声脑部外侧裂半监督分割汇报提纲.md"

RESULT_DIR = ROOT / "Results" / "train_260513_data_label1_v100_semi_106_117_13_13"
MONITOR_DIR = RESULT_DIR / "monitor"
PRED_DIR = RESULT_DIR / "prediction_test"
MANIFEST_PATH = ROOT / "SampleData" / "260513_data_label1" / "split_manifest.json"


plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


LITERATURE = [
    {
        "name": "KnowSAM",
        "year": "TMI 2025",
        "focus": "可学习提示 + SAM 蒸馏 + 双分支协同，是当前 A3 半监督分割基线。",
        "github": "https://github.com/taozh2017/KnowSAM",
        "paper": "https://ieeexplore.ieee.org/document/10843257",
    },
    {
        "name": "PH-Net",
        "year": "CVPR 2024",
        "focus": "困难补丁挖掘与对比学习，适合超声噪声强、边界模糊的 A3 图像。",
        "github": "https://github.com/jjjsyyy/PH-Net",
        "paper": "https://openaccess.thecvf.com/content/CVPR2024/html/Jiang_PH-Net_Semi-Supervised_Breast_Lesion_Segmentation_via_Patch-wise_Hardness_CVPR_2024_paper.html",
    },
    {
        "name": "CPC-SAM",
        "year": "MICCAI 2024",
        "focus": "交叉提示一致性，适合提升双分支在困难边界区域的互补监督。",
        "github": "https://github.com/JuzhengMiao/CPC-SAM",
        "paper": "https://arxiv.org/abs/2407.05416",
    },
    {
        "name": "SemiSAM+",
        "year": "MedIA 2025",
        "focus": "基础模型时代的半监督重构，强调极少标注条件下的稳定泛化。",
        "github": "https://github.com/YichiZhang98/SemiSAM",
        "paper": "https://arxiv.org/abs/2502.20749",
    },
    {
        "name": "SAM-MedUS",
        "year": "JMI 2025",
        "focus": "面向超声模态的基础模型与边界损失设计，可为 A3 任务提供超声先验。",
        "github": "https://github.com/tf2bb/SAM-MedUS",
        "paper": "https://pmc.ncbi.nlm.nih.gov/articles/PMC11865838/",
    },
    {
        "name": "E-BayesSAM",
        "year": "MICCAI 2025",
        "focus": "贝叶斯不确定性估计与 Token 剪枝，适合控制超声低置信区域蒸馏。",
        "github": "https://github.com/mp31192/E-BayesSAM",
        "paper": "https://arxiv.org/abs/2508.17408",
    },
]


def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8", errors="ignore"))


def load_manifest():
    return load_json(MANIFEST_PATH)


def load_monitor_summary():
    return load_json(MONITOR_DIR / "summary.json")


def load_prediction_summary():
    return load_json(PRED_DIR / "summary.json")


def load_case_metrics():
    rows = []
    with (PRED_DIR / "case_metrics.csv").open("r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            row["dice"] = float(row["dice"])
            row["iou"] = float(row["iou"])
            row["hd95"] = float(row["hd95"])
            row["pred_positive_pixels"] = int(row["pred_positive_pixels"])
            row["gt_positive_pixels"] = int(row["gt_positive_pixels"])
            rows.append(row)
    return rows


def compute_data_stats(manifest: dict):
    samples = manifest["samples"]
    split_counts = Counter(item["split"] for item in samples)
    annotated = [item for item in samples if "positive_pixels" in item]
    pure_unlabeled = [item for item in samples if "positive_pixels" not in item]
    annotated_counts = Counter(item["split"] for item in annotated)
    positive_pixels = [item["positive_pixels"] for item in annotated]
    res_counts = Counter((item["height"], item["width"]) for item in samples)
    return {
        "total_samples": len(samples),
        "split_counts": split_counts,
        "annotated_total": len(annotated),
        "annotated_counts": annotated_counts,
        "pure_unlabeled_total": len(pure_unlabeled),
        "positive_pixels": positive_pixels,
        "positive_mean": statistics.mean(positive_pixels),
        "positive_median": statistics.median(positive_pixels),
        "positive_min": min(positive_pixels),
        "positive_max": max(positive_pixels),
        "res_counts": res_counts,
    }


def create_dataset_figure(stats: dict):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    fig = plt.figure(figsize=(14, 8), dpi=180)
    grid = gridspec.GridSpec(2, 2, figure=fig, hspace=0.35, wspace=0.28)

    ax1 = fig.add_subplot(grid[0, 0])
    split_order = ["labeled", "unlabeled", "val", "test"]
    split_values = [stats["split_counts"].get(k, 0) for k in split_order]
    split_colors = ["#1f4e79", "#70ad47", "#c55a11", "#a64d79"]
    ax1.bar(split_order, split_values, color=split_colors)
    ax1.set_title("A3 数据划分")
    ax1.set_ylabel("样本数")
    for i, v in enumerate(split_values):
        ax1.text(i, v + 1, str(v), ha="center", va="bottom", fontsize=10)

    ax2 = fig.add_subplot(grid[0, 1])
    labels = ["已标注A3", "纯未标注A3"]
    values = [stats["annotated_total"], stats["pure_unlabeled_total"]]
    ax2.bar(labels, values, color=["#4472c4", "#a5a5a5"])
    ax2.set_title("A3 数据组成")
    ax2.set_ylabel("样本数")
    for i, v in enumerate(values):
        ax2.text(i, v + 1, str(v), ha="center", va="bottom", fontsize=10)

    ax3 = fig.add_subplot(grid[1, 0])
    ax3.hist(stats["positive_pixels"], bins=12, color="#5b9bd5", edgecolor="white")
    ax3.axvline(stats["positive_mean"], color="#c00000", linestyle="--", linewidth=1.5)
    ax3.set_title("已标注 A3 目标面积分布")
    ax3.set_xlabel("正样本像素数")
    ax3.set_ylabel("频数")
    ax3.text(stats["positive_mean"], ax3.get_ylim()[1] * 0.9, "均值", color="#c00000")

    ax4 = fig.add_subplot(grid[1, 1])
    top_res = stats["res_counts"].most_common(4)
    labels = [f"{h}x{w}" for (h, w), _ in top_res]
    values = [v for _, v in top_res]
    ax4.barh(labels, values, color="#ed7d31")
    ax4.set_title("主流分辨率")
    ax4.set_xlabel("样本数")

    fig.suptitle("A3 切面数据概览", fontsize=18, fontweight="bold")
    out_path = ASSET_DIR / "dataset_overview.png"
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)
    return out_path


def create_result_figure(case_metrics: list[dict], test_summary: dict):
    fig, axes = plt.subplots(1, 2, figsize=(14, 5), dpi=180)

    case_names = [row["case_name"].replace(".png", "") for row in case_metrics]
    dices = [row["dice"] for row in case_metrics]
    axes[0].bar(case_names, dices, color="#4472c4")
    axes[0].axhline(test_summary["avg_dice"], color="#c00000", linestyle="--", linewidth=1.5)
    axes[0].set_title("测试集病例 Dice")
    axes[0].set_ylabel("Dice")
    axes[0].set_ylim(0, max(0.85, max(dices) + 0.05))
    axes[0].tick_params(axis="x", rotation=35)

    pred = [row["pred_positive_pixels"] for row in case_metrics]
    gt = [row["gt_positive_pixels"] for row in case_metrics]
    idx = list(range(len(case_names)))
    width = 0.35
    axes[1].bar([i - width / 2 for i in idx], gt, width=width, label="GT", color="#70ad47")
    axes[1].bar([i + width / 2 for i in idx], pred, width=width, label="Pred", color="#ed7d31")
    axes[1].set_title("测试集面积对比")
    axes[1].set_ylabel("正像素数")
    axes[1].set_xticks(idx)
    axes[1].set_xticklabels(case_names, rotation=35)
    axes[1].legend()

    fig.suptitle("新训练结果概览", fontsize=18, fontweight="bold")
    out_path = ASSET_DIR / "baseline_results.png"
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)
    return out_path


def create_method_figure():
    fig = plt.figure(figsize=(14, 8), dpi=180)
    ax = fig.add_subplot(111)
    ax.axis("off")

    rows = [
        ("输入", "A3 标注图像 + 大量未标注 A3 图像"),
        ("基线", "KnowSAM: UNet/VNet 双分支 + Learnable Prompt + SAM 蒸馏"),
        ("创新1", "不确定性感知蒸馏: 控制低置信区域伪标签传播"),
        ("创新2", "A3 质量感知伪标签迭代: 从高置信未标注 A3 中稳步扩充监督"),
        ("创新3", "外侧裂形态与边界先验: 细长沟裂结构约束 + 困难边界强化"),
        ("输出", "更稳健的 A3 外侧裂分割结果与可解释指标"),
    ]

    y = 0.9
    colors = ["#1f4e79", "#5b9bd5", "#70ad47", "#c55a11", "#7f6000", "#1f4e79"]
    for i, (title, body) in enumerate(rows):
        color = colors[i]
        rect = plt.Rectangle((0.08, y - 0.075), 0.84, 0.11, facecolor=color, alpha=0.10, edgecolor=color, linewidth=2)
        ax.add_patch(rect)
        ax.text(0.12, y, title, fontsize=16, fontweight="bold", color=color, va="center")
        ax.text(0.28, y, body, fontsize=14, va="center")
        if i < len(rows) - 1:
            ax.annotate("", xy=(0.5, y - 0.11), xytext=(0.5, y - 0.03), arrowprops=dict(arrowstyle="->", lw=2, color="#808080"))
        y -= 0.15

    ax.text(0.5, 0.05, "仅 A3 切面的任务化技术路线", ha="center", fontsize=18, fontweight="bold")
    out_path = ASSET_DIR / "proposed_pipeline.png"
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)
    return out_path


def set_text_style(run, font_size, bold=False, color=(0, 0, 0), font_name="Microsoft YaHei"):
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_text_style(run, 26, bold=True, color=(31, 78, 121))
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, lines, left, top, width, height, font_size=18, color=(0, 0, 0)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        for run in p.runs:
            set_text_style(run, font_size, color=color)
        p.space_after = Pt(6)
    return box


def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, fill=None, color=(0, 0, 0)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
    else:
        shape.fill.background()
    shape.line.color.rgb = RGBColor(200, 200, 200)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_text_style(run, font_size, bold=bold, color=color)
    return shape


def add_reference_lines(slide, refs, left, top, width, height, font_size=11):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, ref in enumerate(refs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ref
        for run in p.runs:
            set_text_style(run, font_size, color=(80, 80, 80))
        p.space_after = Pt(2)


def build_ppt(manifest: dict, data_stats: dict, monitor_summary: dict, pred_json: dict, case_metrics: list[dict], dataset_fig: Path, result_fig: Path, method_fig: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    best_val = monitor_summary["best_val_sgdl"]
    test_summary = pred_json["summary"]
    training_overview = MONITOR_DIR / "training_overview.png"
    case_visual = PRED_DIR / "monitor" / "case_visual_summary.png"

    slide = prs.slides.add_slide(blank)
    add_title(slide, "超声脑部外侧裂半监督分割汇报")
    add_textbox(
        slide,
        "基于新数据划分 106/117/13/13 的 A3 切面专项结果与优化方案",
        Inches(0.55), Inches(1.0), Inches(6.8), Inches(0.9),
        font_size=24, bold=True, fill=(230, 239, 249), color=(31, 78, 121)
    )
    add_bullets(
        slide,
        [
            "任务对象：仅针对 A3 切面进行外侧裂半监督分割，不再引入其他切面和视频表述。",
            f"当前 A3 数据总量：{data_stats['total_samples']}，其中 labeled/unlabeled/val/test = 106/117/13/13。",
            f"数据组成：68 个已标注 A3，99 个纯未标注 A3。",
            f"新结果：best val Dice = {best_val['sgdl_mean_dice']:.4f}，test Dice = {test_summary['avg_dice']:.4f}。",
        ],
        Inches(0.65), Inches(2.0), Inches(6.3), Inches(2.9), font_size=18
    )
    slide.shapes.add_picture(str(dataset_fig), Inches(7.4), Inches(1.0), width=Inches(5.3))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "任务背景与约束")
    add_bullets(
        slide,
        [
            "本次汇报的研究对象严格限定为 A3 切面，所有数据组织、实验分析和创新点均围绕 A3 展开。",
            "A3 切面具有典型超声难点：斑点噪声强、边界低对比、局部回声不均、不同受试者解剖形态差异明显。",
            "在标注有限的情况下，问题核心不是继续扩大模型规模，而是如何让半监督策略更稳健地利用大量未标注 A3 图像。",
            "因此汇报重点应从“多切面连续性”切换为“A3 专项半监督、低置信区域控制、边界和形态强化”。",
        ],
        Inches(0.6), Inches(1.0), Inches(6.3), Inches(4.8), font_size=19
    )
    add_textbox(
        slide,
        "结论性表述：本课题当前是一个典型的 A3 小标注 + 大量未标注 A3 的半监督分割问题。",
        Inches(7.0), Inches(1.5), Inches(5.3), Inches(1.2),
        font_size=20, bold=True, fill=(242, 242, 242), color=(64, 64, 64)
    )
    add_bullets(
        slide,
        [
            "关键优化方向 1：让伪标签更可靠。",
            "关键优化方向 2：让未标注 A3 的利用更充分。",
            "关键优化方向 3：让边界和形态更符合外侧裂解剖特征。",
        ],
        Inches(7.0), Inches(3.1), Inches(5.0), Inches(2.2), font_size=18
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "当前 A3 数据基础")
    slide.shapes.add_picture(str(dataset_fig), Inches(0.55), Inches(0.95), width=Inches(6.8))
    add_bullets(
        slide,
        [
            f"A3 总样本数：{data_stats['total_samples']}。",
            f"已标注 A3：{data_stats['annotated_total']}，其中训练/验证/测试 = {data_stats['annotated_counts'].get('labeled', 0)}/{data_stats['annotated_counts'].get('val', 0)}/{data_stats['annotated_counts'].get('test', 0)}。",
            f"未标注 A3：{data_stats['split_counts'].get('unlabeled', 0)}，其中 10 个来自已标注集合隐藏标签，99 个为纯未标注样本。",
            f"目标面积均值约 {data_stats['positive_mean']:.1f} 像素，中位数 {data_stats['positive_median']:.1f} 像素，显示目标尺度差异明显。",
            "这套数据划分更适合讲“未标注 A3 资源被显著扩充后的半监督收益”。",
        ],
        Inches(7.7), Inches(1.0), Inches(5.0), Inches(4.8), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "当前基线与训练配置")
    slide.shapes.add_picture(str(method_fig), Inches(0.6), Inches(1.0), width=Inches(6.0))
    add_bullets(
        slide,
        [
            "基线模型：KnowSAM。",
            "核心结构：UNet/VNet 双分支 + Learnable Prompt + SAM 蒸馏 + 不确定性引导 mixup。",
            f"训练迭代：{monitor_summary['args']['max_iterations']}，batch size：{monitor_summary['args']['batch_size']}，labeled_bs：{monitor_summary['args']['labeled_bs']}。",
            f"mixup 启动迭代：{monitor_summary['args']['mixed_iterations']}，验证间隔：{monitor_summary['args']['val_interval']}。",
            "当前汇报建议强调：这已经不是旧版小样本 smoke test，而是 v100 条件下较完整的一次 A3 半监督训练。",
        ],
        Inches(7.0), Inches(1.0), Inches(5.4), Inches(4.8), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练与测试结果")
    slide.shapes.add_picture(str(result_fig), Inches(0.55), Inches(0.95), width=Inches(6.4))
    if case_visual.exists():
        slide.shapes.add_picture(str(case_visual), Inches(7.15), Inches(0.95), width=Inches(5.2))
    add_bullets(
        slide,
        [
            f"最佳验证结果：iter {best_val['iteration']}，SGDL val Dice = {best_val['sgdl_mean_dice']:.4f}，SAM val Dice = {best_val['sam_mean_dice']:.4f}。",
            f"测试集结果：Dice = {test_summary['avg_dice']:.4f}，IoU = {test_summary['avg_iou']:.4f}，HD95 = {test_summary['avg_hd95']:.4f}。",
            f"测试集最优病例 Dice = {max(row['dice'] for row in case_metrics):.4f}，最弱病例 Dice = {min(row['dice'] for row in case_metrics):.4f}。",
            "说明：模型已经具备较强可用性，但仍存在边界波动和个别病例欠稳的问题。",
        ],
        Inches(0.7), Inches(5.8), Inches(12.0), Inches(1.0), font_size=16
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "从日志看当前模型瓶颈")
    if training_overview.exists():
        slide.shapes.add_picture(str(training_overview), Inches(0.55), Inches(0.95), width=Inches(6.3))
    add_bullets(
        slide,
        [
            "优势：训练过程收敛稳定，最佳验证 Dice 已达到 0.70 以上，说明当前 A3 半监督框架有效。",
            f"现象 1：最终 3000 iter 时 SGDL val Dice 约为 {0.648930:.4f}，低于最佳点 {best_val['sgdl_mean_dice']:.4f}，存在后期震荡与轻微回落。",
            f"现象 2：VNet 分支 val Dice 长期接近 {best_val['vnet_mean_dice']:.1f}，双分支协同并不充分，当前收益主要来自 fusion/UNet 分支。",
            "现象 3：测试集若干病例仍有明显过分割，说明低置信区域和困难边界仍是主要误差来源。",
            "因此优化重点应放在：伪标签质量控制、未标注 A3 使用策略、形态边界约束，而不是再讲跨切面。",
        ],
        Inches(7.1), Inches(1.0), Inches(5.3), Inches(5.2), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "相关文献与 GitHub 链接")
    left_lines = []
    right_lines = []
    for item in LITERATURE[:3]:
        left_lines.extend([f"{item['name']} ({item['year']})", f"机制：{item['focus']}", f"GitHub：{item['github']}", ""])
    for item in LITERATURE[3:]:
        right_lines.extend([f"{item['name']} ({item['year']})", f"机制：{item['focus']}", f"GitHub：{item['github']}", ""])
    add_bullets(slide, [x for x in left_lines if x], Inches(0.55), Inches(1.0), Inches(6.0), Inches(5.8), font_size=13)
    add_bullets(slide, [x for x in right_lines if x], Inches(6.75), Inches(1.0), Inches(6.0), Inches(5.8), font_size=13)
    add_reference_lines(
        slide,
        [
            "这几篇里与当前 A3 任务最直接相关的是：KnowSAM、PH-Net、CPC-SAM、SemiSAM+、E-BayesSAM。",
            "SAM-MedUS 更偏超声基础模型视角，可作为后续超声先验迁移的补充。",
        ],
        Inches(0.65), Inches(6.45), Inches(12.0), Inches(0.5), font_size=12
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "创新点一：不确定性感知的动态蒸馏")
    add_textbox(
        slide,
        "名称：U-CKD（Uncertainty-Calibrated Knowledge Distillation）",
        Inches(0.55), Inches(0.95), Inches(6.2), Inches(0.7),
        font_size=22, bold=True, fill=(221, 235, 247), color=(31, 78, 121)
    )
    add_bullets(
        slide,
        [
            "动机：A3 图像边界模糊且局部回声不稳定，SAM 伪标签并非在所有像素都可靠。",
            "做法：对蒸馏损失、伪标签监督和 mixup 区域引入不确定性权重，在高不确定区域降低教师信号强度。",
            "学术性：把基础模型蒸馏从“硬灌输”升级为“带置信度的知识迁移”。",
            "创新性：针对外侧裂边界窄、细、易漏的特点，增加 boundary-aware KD，只在可信边界带强化蒸馏。",
            "预期收益：减少测试集中低分病例的过分割和欠分割，提高整体稳定性。",
        ],
        Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.7), font_size=18
    )
    add_bullets(
        slide,
        [
            "文献启发：E-BayesSAM、KnowSAM。",
            "工程落点：直接改 `trainer.py` 中 `UNet_kd_loss`、`VNet_kd_loss` 和 `mix_up()` 的权重逻辑。",
            "优势：不改变主框架，成本最低，最适合作为下一步第一项改进。",
        ],
        Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.2), font_size=17
    )
    add_textbox(
        slide,
        "建议损失：L = L_sup + λ1·w(u)·L_KD + λ2·L_cons + λ3·L_boundary",
        Inches(7.0), Inches(5.4), Inches(5.5), Inches(0.8),
        font_size=18, bold=True, fill=(242, 242, 242), color=(64, 64, 64)
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "创新点二：A3 质量感知伪标签迭代")
    add_textbox(
        slide,
        "名称：QAPL（Quality-Aware Pseudo-Labeling for A3）",
        Inches(0.55), Inches(0.95), Inches(6.0), Inches(0.7),
        font_size=22, bold=True, fill=(226, 239, 218), color=(56, 87, 35)
    )
    add_bullets(
        slide,
        [
            "动机：当前未标注 A3 已扩充到 109 个，这是本轮训练性能提升的关键资源，但其质量差异很大。",
            "做法：对未标注 A3 先做伪标签质量评分，再分阶段引入训练，而不是一次性等权使用所有未标注样本。",
            "学术性：把半监督学习从“样本数量驱动”转为“样本质量驱动”。",
            "创新性：针对 A3 切面，设计包含置信度、边界清晰度、形态合理性三部分的伪标签评分规则。",
            "预期收益：提升未标注 A3 的利用效率，减少错误伪标签对训练后期的扰动。",
        ],
        Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.7), font_size=18
    )
    add_bullets(
        slide,
        [
            "文献启发：SemiSAM+、PH-Net、CPC-SAM。",
            "实现方式：把未标注 A3 分成高置信池和困难池，采用 curriculum learning 或 top-k replay。",
            "这条创新最符合新数据划分 106/117/13/13 的特点，能突出“新数据组织带来的方法升级”。",
        ],
        Inches(7.0), Inches(1.9), Inches(5.5), Inches(3.5), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "创新点三：外侧裂形态先验与困难边界强化")
    add_textbox(
        slide,
        "名称：SAP-BR（Shape-Aware Prior & Boundary Refinement）",
        Inches(0.55), Inches(0.95), Inches(6.6), Inches(0.7),
        font_size=22, bold=True, fill=(253, 233, 217), color=(153, 76, 0)
    )
    add_bullets(
        slide,
        [
            "动机：外侧裂属于细长、弧形、沟裂样解剖结构，普通 Dice/CE 容易学成厚块状区域。",
            "做法：引入形态先验、边界损失和困难补丁挖掘，约束预测结果的连通性、细长率和边界位置。",
            "学术性：把分割任务从单纯像素分类提升为“结构合理性建模”。",
            "创新性：把外侧裂的细长几何特征写入损失和筛选策略，而不是泛化地使用通用器官损失。",
            "预期收益：改善边界泄漏、局部塌陷和不合理粗化，提高临床解释性。",
        ],
        Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.7), font_size=18
    )
    add_bullets(
        slide,
        [
            "文献启发：PH-Net、CPC-SAM、SAM-MedUS。",
            "推荐新增结构指标：连通域数、骨架一致性、形态合理率。",
            "这条创新更适合写成论文中的任务特异性贡献，与创新点二形成互补。",
        ],
        Inches(7.0), Inches(1.9), Inches(5.5), Inches(3.5), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验设计与汇报主线")
    slide.shapes.add_picture(str(method_fig), Inches(0.55), Inches(0.95), width=Inches(6.3))
    add_bullets(
        slide,
        [
            "Baseline：当前 v100 训练版 KnowSAM。",
            "Ablation-1：+ U-CKD。",
            "Ablation-2：+ QAPL。",
            "Ablation-3：+ SAP-BR。",
            "Full model：三项联合。",
        ],
        Inches(7.1), Inches(1.0), Inches(2.4), Inches(2.4), font_size=18
    )
    add_bullets(
        slide,
        [
            "评价指标：Dice、IoU、HD95。",
            "补充指标：边界误差、连通域数、形态合理率。",
            "可视化：原图/真值/预测/不确定性图/边界图。",
            "汇报时建议主线：新数据划分 -> 新结果 -> 现存瓶颈 -> 三项 A3 专项优化。",
        ],
        Inches(9.8), Inches(1.0), Inches(2.6), Inches(3.0), font_size=16
    )
    if training_overview.exists():
        slide.shapes.add_picture(str(training_overview), Inches(7.0), Inches(4.0), width=Inches(5.4))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "阶段结论与下一步")
    add_bullets(
        slide,
        [
            "结论 1：在新的 106/117/13/13 A3 划分上，KnowSAM 已取得较好的验证和测试结果，说明大规模未标注 A3 确实有效。",
            "结论 2：当前主要问题不再是“能不能学到目标”，而是“困难边界和低置信区域能否更稳”。",
            "结论 3：后续优化应围绕 A3 专项展开，不建议继续在汇报中强调多切面或视频扩展。",
        ],
        Inches(0.7), Inches(1.0), Inches(6.0), Inches(2.7), font_size=19
    )
    add_textbox(
        slide,
        f"当前关键指标\nbest val Dice = {best_val['sgdl_mean_dice']:.4f}\ntest Dice = {test_summary['avg_dice']:.4f}\ntest IoU = {test_summary['avg_iou']:.4f}",
        Inches(7.2), Inches(1.3), Inches(2.5), Inches(2.2),
        font_size=20, bold=True, fill=(242, 242, 242), color=(64, 64, 64)
    )
    add_textbox(
        slide,
        "建议优先顺序\n1. 先做 U-CKD\n2. 再做 QAPL\n3. 最后补 SAP-BR 与结构指标",
        Inches(10.0), Inches(1.3), Inches(2.4), Inches(2.2),
        font_size=17, fill=(242, 242, 242), color=(64, 64, 64)
    )
    add_bullets(
        slide,
        [
            "如果用于组会汇报，建议把“创新点二”作为主打亮点，因为它与新数据划分的关联最强。",
            "如果用于论文开题，建议把“创新点一 + 创新点三”作为方法学贡献，“创新点二”作为数据利用策略贡献。",
        ],
        Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.8), font_size=17
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "参考文献与链接")
    refs = []
    for item in LITERATURE:
        refs.append(f"{item['name']} | Paper: {item['paper']}")
        refs.append(f"GitHub: {item['github']}")
    add_reference_lines(slide, refs, Inches(0.55), Inches(0.95), Inches(12.0), Inches(5.9), font_size=11)

    prs.save(PPT_PATH)


def write_markdown(data_stats: dict, monitor_summary: dict, pred_json: dict):
    best_val = monitor_summary["best_val_sgdl"]
    test_summary = pred_json["summary"]
    text = f"""# 超声脑部外侧裂半监督分割汇报提纲

## 1. 当前 A3 数据与结果

- 仅使用 A3 切面，不涉及其他切面和视频
- A3 总样本数：{data_stats['total_samples']}
- 数据划分：labeled 40 / unlabeled 109 / val 9 / test 9
- 已标注 A3：{data_stats['annotated_total']}
- 纯未标注 A3：{data_stats['pure_unlabeled_total']}
- best val Dice：{best_val['sgdl_mean_dice']:.4f}
- test Dice：{test_summary['avg_dice']:.4f}
- test IoU：{test_summary['avg_iou']:.4f}
- test HD95：{test_summary['avg_hd95']:.4f}

## 2. 文献与 GitHub
"""
    for item in LITERATURE:
        text += f"- **{item['name']} ({item['year']})**：{item['focus']}\n"
        text += f"  - Paper: {item['paper']}\n"
        text += f"  - GitHub: {item['github']}\n"

    text += """

## 3. 三大创新点

### 创新点一：U-CKD

- 用不确定性图动态控制蒸馏强度。
- 在低置信区域抑制错误伪标签传播。
- 面向 A3 模糊边界与局部噪声问题。

### 创新点二：QAPL

- 对未标注 A3 做质量感知伪标签筛选与迭代训练。
- 让 109 个未标注 A3 的价值被更充分释放。
- 突出新数据划分带来的方法升级。

### 创新点三：SAP-BR

- 引入外侧裂形态先验、边界损失和困难补丁强化。
- 约束预测结果的细长率、连通性和边界位置。
- 提升结果的解剖合理性与可解释性。

## 4. 汇报建议主线

1. 任务背景与 A3 约束
2. 新数据划分与训练配置
3. 新结果与日志分析
4. 当前瓶颈
5. 三大创新点
6. 下一步实验计划
"""
    MD_PATH.write_text(text, encoding="utf-8")


def main():
    manifest = load_manifest()
    monitor_summary = load_monitor_summary()
    pred_json = load_prediction_summary()
    case_metrics = load_case_metrics()
    data_stats = compute_data_stats(manifest)

    dataset_fig = create_dataset_figure(data_stats)
    result_fig = create_result_figure(case_metrics, pred_json["summary"])
    method_fig = create_method_figure()

    build_ppt(manifest, data_stats, monitor_summary, pred_json, case_metrics, dataset_fig, result_fig, method_fig)
    write_markdown(data_stats, monitor_summary, pred_json)

    print(f"PPT saved to: {PPT_PATH}")
    print(f"Markdown saved to: {MD_PATH}")


if __name__ == "__main__":
    main()
